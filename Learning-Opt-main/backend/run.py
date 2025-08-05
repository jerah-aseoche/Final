from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pptx import Presentation
from openpyxl import load_workbook
from datetime import datetime

from app.routes.auth import auth_bp
from app.routes.generate import bp as generate_bp
from app.routes.upload import bp as upload_bp
from app.routes.excel_generate import excel_bp


import os
import uuid
import io, json, re, traceback
import pandas as pd

app = Flask(__name__)
CORS(app)

CORS(app, resources={r"/api/*": {"origins": "*"}})

@app.after_request
def expose_headers(resp):
    resp.headers["Access-Control-Expose-Headers"] = "Content-Disposition"
    return resp

UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), "uploads", "templates")
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

GENERATED_FOLDER = os.path.join("static", "generated")
os.makedirs(GENERATED_FOLDER, exist_ok=True)

# Blueprints
app.register_blueprint(auth_bp)
app.register_blueprint(generate_bp)
app.register_blueprint(upload_bp)
app.register_blueprint(excel_bp)

DEFAULT_MAPPING = {}
PLACEHOLDER_RE = re.compile(r"\{([^}]+)\}")
recent_downloads = []

# quick ping
@app.route("/api/ping")
def ping():
    return jsonify(ok=True)

@app.route("/")
def home():
    return "Hello, Creo Certificate Backend!"

def format_value(val, fmt=None):
    return "" if val is None else str(val)


def replace_placeholders_in_cell(text, mapping, rowdict):
    if "YEAR LAST ATTENDED" in text.upper():
        context = None
        up = text.upper()
        if "ELEMENTARY" in up:
            context = "ELEMENTARY"
        elif "SECONDARY" in up:
            context = "SECONDARY"
        elif "TERTIARY" in up:
            context = "TERTIARY"
    else:
        context = None

    def repl(m):
        key = m.group(1)
        mp = mapping.get(key, key)
        if isinstance(mp, dict):
            col = mp.get(context) or mp.get("DEFAULT")
        else:
            col = mp
        val = rowdict.get(col, "")
        return format_value(val)

    return PLACEHOLDER_RE.sub(repl, text)


def replace_placeholders_in_worksheet(ws, mapping, rowdict):
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        for cell in row:
            if isinstance(cell.value, str) and "{" in cell.value and "}" in cell.value:
                cell.value = replace_placeholders_in_cell(cell.value, mapping, rowdict)


def _safe_sheet_title(s: str, used: set) -> str:
    title = (s or "").strip() or "Row"
    for ch in '[]:*?/\\':
        title = title.replace(ch, "-")
    title = title[:31] or "Row"
    orig = title
    i = 2
    while title in used:
        suffix = f" ({i})"
        title = (orig[: 31 - len(suffix)] + suffix) if len(orig) + len(suffix) > 31 else orig + suffix
        i += 1
    used.add(title)
    return title


def _copy_template_sheet_with_fallback(wb, template_ws, new_title):
    try:
        ws_copy = wb.copy_worksheet(template_ws)
        ws_copy.title = new_title
        return ws_copy
    except Exception as e:
        print("[WARN] copy_worksheet failed; falling back to manual copy:", repr(e))
        ws = wb.create_sheet(title=new_title)
        for rng in template_ws.merged_cells.ranges:
            ws.merge_cells(str(rng))
        for r in range(1, template_ws.max_row + 1):
            for c in range(1, template_ws.max_column + 1):
                v = template_ws.cell(row=r, column=c).value
                if v is not None:
                    ws.cell(row=r, column=c, value=v)
        return ws

@app.route('/generate/certificates', methods=['POST'])
def generate_certificates():
    data = request.json
    template_path = data.get("templatePath")
    output_folder = "static/generated"
    os.makedirs(output_folder, exist_ok=True)

    # Get custom filename from request, or fallback
    custom_filename = data.get("filename")
    if custom_filename:
        filename = f"{custom_filename}.pptx"
    else:
        name = data.get("name", "Certificate")
        filename = f"{name.replace(' ', '_')}_Certificate.pptx"

    output_path = os.path.join(output_folder, filename)

    # Load and customize the PPTX
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{" in run.text and "}}" in run.text:
                            key = run.text.replace("{{", "").replace("}}", "").strip()
                            run.text = data.get(key, "")

    prs.save(output_path)

    # Return list with one file
    return jsonify({"files": [filename]})

@app.route('/api/generate', methods=['POST'])
def generate_tesda_excel():
    uploaded_file = request.files.get("file")
    if not uploaded_file:
        return jsonify({"error": "No file uploaded"}), 400

    # Save temporarily
    temp_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4().hex}.xlsx")
    uploaded_file.save(temp_path)

    try:
        # Load Excel
        wb = load_workbook(temp_path)
        ws = wb.active

        # Save to generated folder
        now = datetime.now().strftime("%Y%m%d-%H%M%S")
        output_filename = f"tesda_record_{now}.xlsx"
        output_path = os.path.join(GENERATED_FOLDER, output_filename)
        wb.save(output_path)

        # ✅ Track in recent_downloads with full metadata for frontend
        recent_downloads.insert(0, {
            "type": "tesda",
            "filename": output_filename,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(output_path)).strftime("%Y-%m-%d %H:%M:%S"),
            "url": f"/static/generated/{output_filename}"
        })

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

def api_generate_certificates():
    import requests

    # 1. Forward request to internal generator
    response = requests.post('http://localhost:5000/generate/certificates', json=request.get_json())
    result = response.json()

    # 2. Handle error if generation failed
    if response.status_code != 200:
        return jsonify({"error": "Failed to generate certificates"}), 500

    # 3. Get list of generated files
    generated_files = result.get("files", [])

    # 4. ✅ Track each generated file in download history
    for fname in generated_files:
        recent_downloads.append({
            "type": "certificate",
            "filename": fname,
        })

    # 5. Return original result to frontend
    return jsonify(result)

@app.route('/api/certificates', methods=['GET'])
def get_certificates():
    try:
        files = [f for f in os.listdir(GENERATED_FOLDER) if f.endswith(".pptx")]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(GENERATED_FOLDER, x)), reverse=True)
        return jsonify(files)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/tesda', methods=['GET'])
def get_tesda_records():
    try:
        files = [f for f in os.listdir(GENERATED_FOLDER) if f.endswith(".xlsx")]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(GENERATED_FOLDER, x)), reverse=True)
        return jsonify(files)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/download-history", methods=["GET"])
def get_download_history():
    
    folder = os.path.join("static", "generated")
    files = [
        f for f in os.listdir(folder)
        if f.endswith(".pptx") or (f.endswith(".xlsx") and "tesda" in f.lower())
        
    ]

    # Sort by last modified time
    files.sort(key=lambda f: os.path.getmtime(os.path.join(folder, f)), reverse=True)

    history = []
    for f in files:
        file_type = "certificate" if f.endswith(".pptx") else "tesda"
        history.append({
            "type": file_type,
            "filename": f,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(os.path.join(folder, f))).strftime("%Y-%m-%d %H:%M"),
            "url": f"/static/generated/{f}"
        })
    return jsonify(history)

@app.route("/api/certificates", methods=["GET"])
def list_certificates():
    folder = os.path.join("static", "generated")
    if not os.path.exists(folder):
        return jsonify([])

    files = [
        f for f in os.listdir(folder)
        if f.endswith(".pptx") and f != "example.pptx"
    ]

    # Sort by last modified time descending
    files.sort(key=lambda f: os.path.getmtime(os.path.join(folder, f)), reverse=True)

    return jsonify(files)

@app.route('/api/tesda')
def list_tesda_files():
    files = [
        f for f in os.listdir("static/generated")
        if f.endswith(".xlsx") and "TESDA" in f
    ]
    return jsonify(files)

# TESDA GENERATION ROUTE (internal)
@app.route('/generate/tesda', methods=['POST'])
def generate_tesda_file():
    try:
        data = request.get_json()
        template_name = data.get("template")
        entries = data.get("data")

        if not template_name or not entries:
            return jsonify({"error": "Missing template or data"}), 400

        template_path = os.path.join("uploads", "templates", template_name)
        if not os.path.exists(template_path):
            return jsonify({"error": "Template not found"}), 404

        base_wb = load_workbook(template_path)
        template_ws = base_wb.active

        used_titles = set()
        for idx, entry in enumerate(entries):
            candidate_name = entry.get("Name", f"Sheet{idx+1}")
            new_title = _safe_sheet_title(candidate_name, used_titles)
            ws_copy = _copy_template_sheet_with_fallback(base_wb, template_ws, new_title)
            replace_placeholders_in_worksheet(ws_copy, {}, entry)

        base_wb.remove(template_ws)

        filename = f"TESDA_{datetime.now().strftime('%Y-%m-%d_%H:%M')}.xlsx"
        output_path = os.path.join("static", "generated", filename)
        base_wb.save(output_path)

        return jsonify({"files": [filename]}), 200

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500



# API ROUTE THAT CALLS INTERNAL GENERATOR AND TRACKS HISTORY
# ✅ TESDA generation route (calls internal generator and logs to history)
@app.route('/api/generate-tesda', methods=['POST'])
def api_generate_tesda():
    # Forward request to the internal generator endpoint
    response = requests.post('http://localhost:5000/generate/tesda', json=request.get_json())

    if response.status_code != 200:
        return jsonify({"error": "Failed to generate TESDA file"}), 500

    result = response.json()
    generated_files = result.get("files", [])

    for fname in generated_files:
        file_path = os.path.join("static", "generated", fname)
        if os.path.exists(file_path):  # ✅ Only add to history if file exists
            recent_downloads.insert(0, {
                "type": "tesda",
                "filename": fname,
                "timestamp": datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M:%S"),
                "url": f"/static/generated/{fname}"
            })

    return jsonify(result)


# ✅ Used by frontend to track downloads and update history
@app.route("/api/download-history", methods=["POST"])
def update_download_history():
    data = request.get_json()
    filename = data.get("filename")
    if not filename:
        return jsonify({"error": "Missing filename"}), 400

    file_path = os.path.join("static", "generated", filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File does not exist"}), 404

    # Avoid duplicates
    if not any(d.get("filename") == filename for d in recent_downloads):
        file_type = "tesda" if filename.lower().endswith(".xlsx") else "certificate"
        recent_downloads.insert(0, {
            "type": file_type,
            "filename": filename,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M"),
            "url": f"/static/generated/{filename}"
        })

    return jsonify({"success": True})

if __name__ == "__main__":
    app.run(debug=True, host='0.0.0.0', port=5000)
