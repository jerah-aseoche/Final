from flask import Blueprint, request, jsonify, send_file
from flask_cors import cross_origin
from pptx import Presentation
import os
import json
import io
import re
from copy import deepcopy
from datetime import datetime


bp = Blueprint('generate', __name__, url_prefix='/generate')

# Paths
BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))
TEMPLATE_DIR = os.path.join(BASE_DIR, 'uploads', 'templates')
OUTPUT_DIR = os.path.join(BASE_DIR, 'static', 'generated')

os.makedirs(OUTPUT_DIR, exist_ok=True)

# Helper to fill slide placeholders
def fill_slide(slide, data_row):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = ''.join(run.text for run in paragraph.runs)
            replaced = full_text
            for key, value in data_row.items():
                ph = f"{{{key}}}"
                replaced = replaced.replace(ph, str(value))
            if replaced != full_text:
                for run in paragraph.runs:
                    run.text = ""
                paragraph.runs[0].text = replaced

@bp.route('/delete_certificate', methods=['DELETE'])
@cross_origin()
def delete_certificate():
    filename = request.args.get("filename")
    if not filename:
        return jsonify({"error": "Filename is required"}), 400

    file_path = os.path.join(OUTPUT_DIR, filename)

    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404

    try:
        os.remove(file_path)
        return jsonify({"message": "File deleted successfully"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@bp.route('/certificates', methods=['POST', 'OPTIONS'])
@cross_origin()
def generate_certificates():
    data = request.get_json()
    template_type = data.get('template', 'ojt')
    rows = data.get('rows', [])

    if not rows:
        return jsonify({"error": "No data provided"}), 400

    tpl_filename = f"{template_type}.pptx"
    template_path = os.path.join(TEMPLATE_DIR, tpl_filename)
    if not os.path.exists(template_path):
        return jsonify({"error": f"Template '{tpl_filename}' not found"}), 404

    prs = Presentation(template_path)
    source_slide = prs.slides[0]
    original_elements = [deepcopy(shape.element) for shape in source_slide.shapes]
    fill_slide(source_slide, rows[0])

    for row in rows[1:]:
        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        for shp in list(new_slide.shapes):
            new_slide.shapes._spTree.remove(shp.element)
        for el in original_elements:
            new_slide.shapes._spTree.append(deepcopy(el))
        fill_slide(new_slide, row)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H:%M")
    output_name = f"certificate_{template_type} ({timestamp}).pptx"
    output_path = os.path.join(OUTPUT_DIR, output_name)
    prs.save(output_path)

    return jsonify({"message": "Certificates generated", "files": [output_name]})

@bp.route('/files/<filename>', methods=['GET'])
@cross_origin()
def get_generated_file(filename):
    file_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404
    return send_file(file_path, as_attachment=True, download_name=filename)  # ✅ Fix here

@bp.route('/preview', methods=['POST', 'OPTIONS'])
@cross_origin()
def preview_certificate():
    data = request.get_json()
    template_type = data.get('template', 'ojt')
    rows = data.get('rows', [])

    if not rows:
        return jsonify({"error": "No data to preview"}), 400

     
    tpl_filename = f"{template_type}.pptx"
    template_path = os.path.join(TEMPLATE_DIR, tpl_filename)

    if not os.path.exists(template_path):
        return jsonify({"error": f"Template '{tpl_filename}' not found"}), 404

    if template_type == 'tesda':
        def get_value(row, target_key):
            for k in row.keys():
                if k.strip().lower() == target_key.strip().lower():
                    return row[k]
            return ""

        def field_row(label, value):
            return f"<div class='tesda-field-row'><div class='tesda-label'>{label}</div><div class='tesda-value'>{value}</div></div>"

        html_parts = [
            "<!DOCTYPE html>",
            "<html><head><meta charset='utf-8'>",
            "<meta name='viewport' content='width=device-width, initial-scale=1.0'>",
            "<title>TESDA Data Preview</title><style>",
            "body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: #2d3748; margin: 0; padding: 0; display: flex; justify-content: center; flex-direction: column; }",
            ".container { max-width: 1300px; width: 100%; padding: 10px; text-align: center; display: flex; flex-direction: column; height: 100vh; justify-content: space-between; }",
            ".top-nav { display: flex; justify-content: space-between; align-items: center; margin-bottom: 20px; }",
            ".top-nav h2 { color: white; margin: 0; font-size: 1.8rem; }",
            ".top-nav button { background-color: #a361ef; border: none; padding: 8px 16px; border-radius: 8px; color: white; cursor: pointer; font-weight: bold; }",
            ".top-nav button:disabled { background-color: #888; cursor: not-allowed; }",
            ".tesda-card { max-width: 1000px; background: transparent; padding: 20px; display: none; flex-direction: column; align-items: center; margin: auto; }",
            ".tesda-card.active { display: flex; }",
            ".tesda-pair { display: flex; gap: 30px; justify-content: center; width: 100%; }",
            ".tesda-profile { background: #4a5568; border-radius: 16px; box-shadow: 0 4px 24px rgba(0,0,0,0.18); padding: 10px 10px; flex: 1; min-width: 600px; }",
            ".tesda-profile.empty { background: transparent; box-shadow: none; }",
            ".tesda-profile h3 { color: #fff; font-size: 1.5rem; margin-bottom: 16px; text-align: center; }",
            ".tesda-field-row { display: flex; justify-content: space-between; width: 100%; margin-bottom: 12px; }",
            ".tesda-label { color: #000; font-weight: 600; width: 200px; text-align: left; }",
            ".tesda-value { color: #e2e8f0; flex: 1; text-align: left; word-break: break-word; }",
            "</style></head><body><div class='container'>",

            # 👇 Top navigation bar with header + buttons
            "<div class='top-nav'>",
            "<button id='prevBtn' onclick='showCard(currentIndex - 1)' disabled>Previous</button>",
            "<h2>TESDA Data Preview</h2>",
            "<button id='nextBtn' onclick='showCard(currentIndex + 1)'>Next</button>",
            "</div>"
        ]

        for index in range(0, len(rows), 2):
            row1 = rows[index]
            row2 = rows[index + 1] if index + 1 < len(rows) else None

            html_parts.append(f"<div class='tesda-card{' active' if index == 0 else ''}' id='card-{index // 2}'>")
            html_parts.append("<div class='tesda-pair'>")

            def render_profile(row):
                if not row:
                    return "<div class='tesda-profile empty'></div>"
                name = get_value(row, "NAME") or "No Name"
                profile_html = f"<div class='tesda-profile'><h3>{name}</h3>"
                profile_html += "".join([
                    field_row("Gender:", get_value(row, "Gender")),
                    field_row("Date of Birth:", get_value(row, "Date of Birth")),
                    field_row("Place of Birth:", get_value(row, "Place of Birth")),
                    field_row("Home Address:", get_value(row, "Home Address")),
                    field_row("Elementary:", get_value(row, "Elementary")),
                    field_row("Year Attended:", get_value(row, "Year Last Attended - Elementary")),
                    field_row("Secondary:", get_value(row, "Secondary")),
                    field_row("Year Attended:", get_value(row, "Year Last Attended - Secondary")),
                    field_row("Tertiary:", get_value(row, "Tertiary")),
                    field_row("Year Attended:", get_value(row, "Year Last Attended - Tertiary")),
                    field_row("Date Started:", get_value(row, "Date Started")),
                    field_row("Date Finished:", get_value(row, "Date Finished")),
                    field_row("Date of Graduation:", get_value(row, "Date of Graduation")),
                ])
                profile_html += "</div>"
                return profile_html

            html_parts.append(render_profile(row1))
            html_parts.append(render_profile(row2))
            html_parts.append("</div></div>")  # Close tesda-pair and tesda-card

        total_pages = (len(rows) + 1) // 2

        html_parts.append(f"""
            </div> <!-- close container -->
            <script>
                let currentIndex = 0;
                const total = {total_pages};
                function showCard(index) {{
                    if (index < 0 || index >= total) return;
                    document.getElementById('card-' + currentIndex).classList.remove('active');
                    currentIndex = index;
                    document.getElementById('card-' + currentIndex).classList.add('active');
                    document.getElementById('prevBtn').disabled = currentIndex === 0;
                    document.getElementById('nextBtn').disabled = currentIndex === total - 1;
                }}
            </script>
            </body></html>
        """)

        return "\n".join(html_parts), 200, {"Content-Type": "text/html"}

    # Default preview (non-TESDA)
    html_parts = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'>",
        "<meta name='viewport' content='width=device-width, initial-scale=1.0'>",
        "<title>Certificate Preview</title><style>",
        "body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: #2d3748; margin: 0; padding: 20px; min-height: 100vh; }",
        ".container { max-width: 900px; margin: 0 auto; }",
        "h2 { text-align: center; color: white; font-size: 2.5rem; margin-bottom: 30px; }",
        ".slide-preview { background: #4a5568; border-radius: 12px; padding: 25px; margin-bottom: 25px; box-shadow: 0 8px 32px rgba(0,0,0,0.2); border: 1px solid #718096; transition: transform 0.2s ease; }",
        ".slide-preview:hover { transform: translateY(-2px); box-shadow: 0 12px 40px rgba(0,0,0,0.3); }",
        ".slide-preview h4 { color: white; font-size: 1.3rem; margin: 0 0 15px 0; padding-bottom: 8px; border-bottom: 2px solid #718096; }",
        ".slide-preview p { margin: 10px 0; font-size: 1.1rem; line-height: 1.6; color: #e2e8f0; padding: 8px 12px; background: #2d3748; border-radius: 6px; border-left: 4px solid #a361ef; }",
        ".certificate-text { font-weight: 500; }",
        "@media (max-width: 768px) { .container { padding: 10px; } .slide-preview { padding: 20px; } h2 { font-size: 2rem; } }",
        "</style></head><body><div class='container'><h2>Certificate Preview</h2>"
    ]

    for idx, row in enumerate(rows):
        prs_row = Presentation(template_path)
        slide = prs_row.slides[0]
        fill_slide(slide, row)
        html_parts.append("<div class='slide-preview certificate-text'>")
        html_parts.append(f"<h4>Certificate {idx+1}</h4>")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in paragraph.runs)
                if text.strip():
                    html_parts.append(f"<p>{text}</p>")
        html_parts.append("</div>")
    html_parts.append("</div></body></html>")
    return "\n".join(html_parts), 200, {"Content-Type": "text/html"}
